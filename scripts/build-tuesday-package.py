#!/usr/bin/env python3
"""Build the Tuesday package for Andrew Greene.

Produces:
  docs/proposals/defense-unicorns/tuesday/a11oy_uds_vision_deck.pptx
  docs/proposals/defense-unicorns/tuesday/a11oy_uds_vision_deck.pdf
  docs/proposals/defense-unicorns/tuesday/email_to_andrew.docx
  docs/proposals/defense-unicorns/tuesday/a11oy_uds_package.zip

Run: python3 scripts/build-tuesday-package.py
"""
from __future__ import annotations

import os
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGB, Inches as DocxInches

from reportlab.lib.pagesizes import landscape
from reportlab.pdfgen import canvas as pdf_canvas
from reportlab.lib.colors import HexColor
from reportlab.lib.utils import simpleSplit


ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs/proposals/defense-unicorns/tuesday"

# UDS-inspired palette (dark + purple) with a11oy gold accents
BG = (10, 10, 10)
CARD = (20, 20, 20)
TEXT = (245, 245, 245)
SUB = (138, 138, 138)
GOLD = (201, 183, 135)
UDS = (125, 76, 255)
OK = (74, 222, 128)


# --------------------------------------------------------------------------
# Slide content
# --------------------------------------------------------------------------
SLIDES = [
    {
        "kind": "title",
        "eyebrow": "FOR ANDREW GREENE · DEFENSE UNICORNS · TUESDAY 2026-05-19",
        "title": "a11oy.UDS",
        "subtitle": "A UDS-native, governed agent runtime.\nInheriting UDS guardrails. Carrying a11oy's orchestration DNA.",
        "footer": "Lutar, Stephen P. · ORCID 0009-0001-0110-4173 · SZL Holdings",
    },
    {
        "kind": "section",
        "eyebrow": "POSTURE · SLIDE 02",
        "title": "Why this name, why this shape",
        "body": [
            "The Linux ethos thread: the kernel surface stays small, sharp, and verifiable.",
            "Capabilities mesh in at the policy and admission edges.",
            "",
            "a11oy is a capability. UDS is the kernel surface.",
        ],
    },
    {
        "kind": "grid",
        "eyebrow": "INHERITANCE · SLIDE 03",
        "title": "What a11oy already carries",
        "items": [
            ("1", "Orchestration plane", "signal mesh, planner, workcells"),
            ("2", "Approval gates", "no material action without human OK"),
            ("3", "Artifact registry", "models, prompts, embeddings, evals, agents"),
            ("4", "Proof ledger", "append-only, Ed25519 + ML-DSA-65"),
            ("5", "Λ-9 invariant runtime", "Doctrine V6 — 0.90 / 0.95 / 0.95"),
            ("6", "Recalibration memo", "weekly 'what changed' feed"),
        ],
    },
    {
        "kind": "grid",
        "eyebrow": "INHERITANCE · SLIDE 04",
        "title": "What UDS already carries",
        "items": [
            ("01", "Distribution", "Zarf packages, UDS bundles"),
            ("02", "Cluster", "uds-core, Pepr admission, NetworkPolicies"),
            ("03", "Identity", "Keycloak SSO, tenant realms"),
            ("04", "Edge", "Istio tenant gateway"),
            ("05", "Telemetry", "Loki, Prometheus, Grafana"),
            ("06", "Distribution at scale", "OCI registry + SBOM flow"),
        ],
    },
    {
        "kind": "table",
        "eyebrow": "THESIS · SLIDE 05",
        "title": "The meshing thesis",
        "headers": ["a11oy primitive", "UDS primitive it inherits from"],
        "rows": [
            ["Orchestration", "Pepr operators, Istio tenant gateway"],
            ["Approval gates", "UDS policy engine, Pepr admission"],
            ["Artifact registry", "Zarf packages, OCI registry, SBOM flow"],
            ["Proof ledger", "Loki + signed attestation sidecar"],
            ["Λ-9 invariant", "Pepr admission module (#5027, merged)"],
            ["Recalibration memo", "NetworkPolicy-aware cluster inventory feed"],
        ],
    },
    {
        "kind": "section",
        "eyebrow": "PROBLEM 1 · SLIDE 06",
        "title": "Trusted AI/agent orchestration\ninside air-gapped UDS",
        "body": [
            "Strong image / chart provenance.   ✓",
            "Strong identity (Keycloak).        ✓",
            "Strong policy (Pepr).              ✓",
            "",
            "Governed agent runtime with structural approval gates,",
            "immutable tool-call audit, and verifiable offline.   ← gap",
        ],
        "accent": UDS,
    },
    {
        "kind": "section",
        "eyebrow": "PROBLEM 2 · SLIDE 07",
        "title": "A UDS-native artifact spine for AI",
        "body": [
            "Container images: signed, versioned, attested, lifecycled.   ✓",
            "",
            "Models, prompts, embeddings, agent definitions, evals:",
            "  SBOM-style attestation per artifact",
            "  Signed evals (dev pass means prod pass)",
            "  Drift detection",
            "  Promote / queue / discard, mirroring Zarf",
        ],
        "accent": GOLD,
    },
    {
        "kind": "ladder",
        "eyebrow": "THE LADDER · SLIDE 08",
        "title": "A / B / C — one ladder, one direction",
        "options": [
            ("OPTION A", "2–3 week proof point", "a11oy.UDS as a Zarf bundle payload. Drops into an existing UDS cluster. Inherits everything as-is.", GOLD),
            ("OPTION B", "Falls out of A", "Primitives ported one-by-one to native UDS components. Done as a by-product of the bundle work.", SUB),
            ("OPTION C", "The real destination", "Full a11oy.UDS ecosystem port. First-class peer of uds-core. Default checkbox at install.", UDS),
        ],
    },
    {
        "kind": "weeks",
        "eyebrow": "PROOF PLAN · SLIDE 09",
        "title": "The 2–3 week proof plan (Option A)",
        "weeks": [
            ("WEEK 1", "Bundle, identity, observability", [
                "uds-cli bundle deploy on the reference cluster",
                "Keycloak SSO round-trip end-to-end",
                "Istio tenant gateway routes verified",
                "Loki + Prometheus exporters confirmed",
            ]),
            ("WEEK 2", "Approval gates, Λ-9 admission, audit chain", [
                "Bad invocation → MATURITY_GATE_BLOCKED",
                "Good invocation → approval queue prompt",
                "Cable pulled — three offline invocations",
                "verify --offline → OK chain=clean",
            ]),
            ("WEEK 3", "Artifact spine", [
                "AIArtifact CRD applied",
                "All five kinds round-trip lifecycle",
                "Seeded drift surfaces in next memo",
                "Broken-signature promote denied + recorded",
            ]),
        ],
    },
    {
        "kind": "section",
        "eyebrow": "CREDIBILITY · SLIDE 10",
        "title": "\"The wires are set up\"",
        "body": [
            "MERGED:",
            "  uds-cli #5026 — in-bundle hash-chained attestation manifest",
            "  pepr #5027  — Λ-floor admission module (0.90 / 0.95 / 0.95)",
            "  #5028       — three Zarf packages + UDS bundle, attestations sidecar wired",
            "",
            "LIVE:",
            "  OPA gateway pack — 3 tests, pinned OPA v0.69.0",
            "",
            "TRACKED:",
            "  #5118 / #5119 — publish + validate housekeeping (out of scope for Tuesday)",
        ],
        "accent": OK,
    },
    {
        "kind": "section",
        "eyebrow": "SO WHAT · SLIDE 11",
        "title": "What credibly-auditable agents unlock\nfor the DoD side",
        "body": [
            "Defensible deployment of LLM-driven decision support inside SCIFs.",
            "Cross-coalition evidence sharing without releasing the model.",
            "Insurance-grade posture claims for AI-driven operations.",
            "A foundation for autonomy the program manager can sign off on.",
        ],
    },
    {
        "kind": "section",
        "eyebrow": "RISKS · SLIDE 12",
        "title": "Risks we've named",
        "body": [
            "License surface — AGPL ↔ Apache managed by dual-licensed upstream PRs.",
            "Doctrine drift — Λ floors payload-anchored; any change is a replay event.",
            "Scope creep — Option C deliberately gated on Andrew's response to A.",
        ],
    },
    {
        "kind": "section",
        "eyebrow": "THE ASK · SLIDE 13",
        "title": "Four small asks",
        "body": [
            "1.  A Mission App target for the Week-3 demo (bland-but-real is best).",
            "2.  Keycloak realm access on a reference cluster.",
            "3.  A 30-minute review window at the end of Week 3.",
            "4.  A thumbs-up to schedule the Option C scoping conversation.",
        ],
        "accent": GOLD,
    },
    {
        "kind": "title",
        "eyebrow": "CLOSE · SLIDE 14",
        "title": "Governed agency",
        "subtitle": "the next primitive the UDS surface deserves —\nand we can prove it in three weeks without you taking our word for any of it.",
        "footer": "Lutar, Stephen P. · ORCID 0009-0001-0110-4173 · SZL Holdings",
        "accent": UDS,
    },
]


# --------------------------------------------------------------------------
# PPTX builder
# --------------------------------------------------------------------------
def _rgb(triple):
    from pptx.dml.color import RGBColor as RGB
    return RGB(*triple)


def _solid_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(color)
    return bg


def _accent_bar(slide, color, x=Inches(0.5), y=Inches(0.55), w=Inches(0.06), h=Inches(0.35)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(color)
    return bar


def _text(slide, text, x, y, w, h, *, size=18, color=TEXT, bold=False, mono=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        if mono:
            run.font.name = "JetBrains Mono"
        else:
            run.font.name = "Inter"
    return box


def _slide_chrome(prs, accent=UDS):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _solid_bg(s, BG)
    # Top accent strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.06))
    strip.line.fill.background()
    strip.fill.solid()
    strip.fill.fore_color.rgb = _rgb(accent)
    # Footer
    _text(s, "a11oy.UDS  ·  for Andrew Greene  ·  Tuesday 2026-05-19",
          Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
          size=9, color=SUB, mono=True)
    return s


def _slide_title(prs, slide):
    s = _slide_chrome(prs, slide.get("accent", UDS))
    _text(s, slide["eyebrow"], Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.3),
          size=10, color=slide.get("accent", UDS), bold=True, mono=True)
    _text(s, slide["title"], Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.0),
          size=84, color=TEXT, bold=True)
    _text(s, slide["subtitle"], Inches(0.5), Inches(4.4), Inches(12.3), Inches(1.6),
          size=18, color=SUB)
    _text(s, slide["footer"], Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.3),
          size=10, color=SUB, mono=True)


def _slide_section(prs, slide):
    accent = slide.get("accent", UDS)
    s = _slide_chrome(prs, accent)
    _accent_bar(s, accent)
    _text(s, slide["eyebrow"], Inches(0.7), Inches(0.55), Inches(12.0), Inches(0.3),
          size=10, color=accent, bold=True, mono=True)
    _text(s, slide["title"], Inches(0.7), Inches(1.0), Inches(12.0), Inches(1.6),
          size=40, color=TEXT, bold=True)
    body = "\n".join(slide["body"])
    _text(s, body, Inches(0.7), Inches(3.0), Inches(12.0), Inches(3.5),
          size=18, color=SUB, mono=any("→" in l or "✓" in l for l in slide["body"]))


def _slide_grid(prs, slide):
    accent = slide.get("accent", UDS)
    s = _slide_chrome(prs, accent)
    _accent_bar(s, accent)
    _text(s, slide["eyebrow"], Inches(0.7), Inches(0.55), Inches(12.0), Inches(0.3),
          size=10, color=accent, bold=True, mono=True)
    _text(s, slide["title"], Inches(0.7), Inches(1.0), Inches(12.0), Inches(0.9),
          size=36, color=TEXT, bold=True)

    items = slide["items"]
    cols = 3
    rows = (len(items) + cols - 1) // cols
    cw, ch = Inches(4.0), Inches(1.7)
    gx, gy = Inches(0.7), Inches(2.4)
    gap_x, gap_y = Inches(0.1), Inches(0.15)
    for i, (n, name, desc) in enumerate(items):
        r, c = divmod(i, cols)
        x = Emu(int(gx) + c * (int(cw) + int(gap_x)))
        y = Emu(int(gy) + r * (int(ch) + int(gap_y)))
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, ch)
        card.line.color.rgb = _rgb((40, 40, 40))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(CARD)
        _text(s, n, Emu(int(x) + Inches(0.2)), Emu(int(y) + Inches(0.15)),
              Inches(1.0), Inches(0.4), size=14, color=accent, bold=True, mono=True)
        _text(s, name, Emu(int(x) + Inches(0.2)), Emu(int(y) + Inches(0.55)),
              Inches(3.7), Inches(0.45), size=15, color=TEXT, bold=True)
        _text(s, desc, Emu(int(x) + Inches(0.2)), Emu(int(y) + Inches(1.0)),
              Inches(3.7), Inches(0.6), size=11, color=SUB)


def _slide_table(prs, slide):
    accent = slide.get("accent", UDS)
    s = _slide_chrome(prs, accent)
    _accent_bar(s, accent)
    _text(s, slide["eyebrow"], Inches(0.7), Inches(0.55), Inches(12.0), Inches(0.3),
          size=10, color=accent, bold=True, mono=True)
    _text(s, slide["title"], Inches(0.7), Inches(1.0), Inches(12.0), Inches(0.9),
          size=36, color=TEXT, bold=True)
    rows = slide["rows"]
    headers = slide["headers"]
    n = len(rows) + 1
    table_shape = s.shapes.add_table(n, len(headers), Inches(0.7), Inches(2.3),
                                     Inches(12.0), Inches(0.55 * n))
    tbl = table_shape.table
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(11)
                r.font.bold = True
                r.font.color.rgb = _rgb(SUB)
                r.font.name = "JetBrains Mono"
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(BG)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(13)
                    r.font.color.rgb = _rgb(TEXT if ci == 0 else SUB)
                    r.font.name = "Inter"
                    if ci == 0:
                        r.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(CARD)


def _slide_ladder(prs, slide):
    s = _slide_chrome(prs, UDS)
    _accent_bar(s, UDS)
    _text(s, slide["eyebrow"], Inches(0.7), Inches(0.55), Inches(12.0), Inches(0.3),
          size=10, color=UDS, bold=True, mono=True)
    _text(s, slide["title"], Inches(0.7), Inches(1.0), Inches(12.0), Inches(0.9),
          size=36, color=TEXT, bold=True)
    cw, ch = Inches(4.0), Inches(3.6)
    gx, gy = Inches(0.7), Inches(2.4)
    gap_x = Inches(0.15)
    for i, (tag, badge, body, color) in enumerate(slide["options"]):
        x = Emu(int(gx) + i * (int(cw) + int(gap_x)))
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, gy, cw, ch)
        card.line.color.rgb = _rgb((40, 40, 40))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(CARD)
        # left accent
        accent_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, gy, Inches(0.06), ch)
        accent_bar.line.fill.background()
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = _rgb(color)
        _text(s, tag, Emu(int(x) + Inches(0.25)), Emu(int(gy) + Inches(0.2)),
              Inches(3.6), Inches(0.4), size=12, color=color, bold=True, mono=True)
        _text(s, badge, Emu(int(x) + Inches(0.25)), Emu(int(gy) + Inches(0.65)),
              Inches(3.6), Inches(0.4), size=11, color=SUB, mono=True)
        _text(s, body, Emu(int(x) + Inches(0.25)), Emu(int(gy) + Inches(1.4)),
              Inches(3.55), Inches(2.0), size=14, color=TEXT)


def _slide_weeks(prs, slide):
    s = _slide_chrome(prs, GOLD)
    _accent_bar(s, GOLD)
    _text(s, slide["eyebrow"], Inches(0.7), Inches(0.55), Inches(12.0), Inches(0.3),
          size=10, color=GOLD, bold=True, mono=True)
    _text(s, slide["title"], Inches(0.7), Inches(1.0), Inches(12.0), Inches(0.9),
          size=36, color=TEXT, bold=True)
    cw, ch = Inches(4.0), Inches(4.3)
    gx, gy = Inches(0.7), Inches(2.3)
    gap_x = Inches(0.15)
    for i, (tag, name, items) in enumerate(slide["weeks"]):
        x = Emu(int(gx) + i * (int(cw) + int(gap_x)))
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, gy, cw, ch)
        card.line.color.rgb = _rgb((40, 40, 40))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(CARD)
        accent_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, gy, cw, Inches(0.05))
        accent_bar.line.fill.background()
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = _rgb(GOLD)
        _text(s, tag, Emu(int(x) + Inches(0.25)), Emu(int(gy) + Inches(0.25)),
              Inches(3.6), Inches(0.4), size=12, color=GOLD, bold=True, mono=True)
        _text(s, name, Emu(int(x) + Inches(0.25)), Emu(int(gy) + Inches(0.7)),
              Inches(3.6), Inches(0.8), size=15, color=TEXT, bold=True)
        body = "\n".join(f"• {it}" for it in items)
        _text(s, body, Emu(int(x) + Inches(0.25)), Emu(int(gy) + Inches(1.8)),
              Inches(3.55), Inches(2.3), size=12, color=SUB)


def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    builders = {
        "title": _slide_title,
        "section": _slide_section,
        "grid": _slide_grid,
        "table": _slide_table,
        "ladder": _slide_ladder,
        "weeks": _slide_weeks,
    }
    for slide in SLIDES:
        builders[slide["kind"]](prs, slide)
    prs.save(path)


# --------------------------------------------------------------------------
# PDF builder (matches deck flow, lightweight, landscape)
# --------------------------------------------------------------------------
def _hex(t):
    return HexColor("#%02x%02x%02x" % t)


def _pdf_chrome(c, w, h, accent):
    c.setFillColor(_hex(BG))
    c.rect(0, 0, w, h, fill=1, stroke=0)
    c.setFillColor(_hex(accent))
    c.rect(0, h - 6, w, 6, fill=1, stroke=0)
    c.setFillColor(_hex(SUB))
    c.setFont("Courier", 8)
    c.drawString(36, 24, "a11oy.UDS  ·  for Andrew Greene  ·  Tuesday 2026-05-19")


def _pdf_wrap(c, text, x, y, max_w, size, font="Helvetica", line_h=None):
    line_h = line_h or size * 1.4
    for raw_line in text.split("\n"):
        if not raw_line.strip():
            y -= line_h
            continue
        for chunk in simpleSplit(raw_line, font, size, max_w):
            c.drawString(x, y, chunk)
            y -= line_h
    return y


def build_pdf(path: Path) -> None:
    W, H = landscape((1280, 720))
    c = pdf_canvas.Canvas(str(path), pagesize=(W, H))
    c.setAuthor("Lutar, Stephen P. (SZL Holdings)")
    c.setCreator("Lutar, Stephen P. (SZL Holdings)")
    c.setTitle("a11oy.UDS - vision deck for Andrew Greene")
    c.setSubject("a11oy.UDS Tuesday package")
    c.setKeywords("a11oy UDS Defense Unicorns SZL Holdings")
    for slide in SLIDES:
        accent = slide.get("accent", UDS)
        _pdf_chrome(c, W, H, accent)
        eyebrow = slide.get("eyebrow", "")
        c.setFillColor(_hex(accent))
        c.setFont("Courier-Bold", 9)
        c.drawString(36, H - 50, eyebrow)
        c.setFillColor(_hex(TEXT))

        kind = slide["kind"]
        if kind == "title":
            c.setFont("Helvetica-Bold", 64)
            for i, line in enumerate(slide["title"].split("\n")):
                c.drawString(36, H - 200 - i * 70, line)
            c.setFillColor(_hex(SUB))
            c.setFont("Helvetica", 16)
            _pdf_wrap(c, slide["subtitle"], 36, H - 360, W - 72, 16)
            c.setFillColor(_hex(SUB))
            c.setFont("Courier", 9)
            c.drawString(36, 60, slide.get("footer", ""))
        elif kind == "section":
            c.setFillColor(_hex(TEXT))
            c.setFont("Helvetica-Bold", 36)
            for i, line in enumerate(slide["title"].split("\n")):
                c.drawString(36, H - 110 - i * 44, line)
            c.setFillColor(_hex(SUB))
            c.setFont("Helvetica", 16)
            y = H - 240
            for line in slide["body"]:
                if any(s in line for s in ("→", "✓", "MERGED", "LIVE", "TRACKED")):
                    c.setFont("Courier", 14)
                else:
                    c.setFont("Helvetica", 16)
                c.drawString(36, y, line)
                y -= 26
        elif kind == "grid":
            c.setFillColor(_hex(TEXT))
            c.setFont("Helvetica-Bold", 32)
            c.drawString(36, H - 100, slide["title"])
            cols = 3
            cw = (W - 72 - 20) / cols
            ch = 130
            for i, (n, name, desc) in enumerate(slide["items"]):
                r, col = divmod(i, cols)
                x = 36 + col * (cw + 10)
                y = H - 180 - r * (ch + 12)
                c.setFillColor(_hex(CARD))
                c.rect(x, y - ch, cw, ch, fill=1, stroke=0)
                c.setFillColor(_hex(accent))
                c.setFont("Courier-Bold", 12)
                c.drawString(x + 14, y - 28, n)
                c.setFillColor(_hex(TEXT))
                c.setFont("Helvetica-Bold", 14)
                c.drawString(x + 14, y - 56, name)
                c.setFillColor(_hex(SUB))
                c.setFont("Helvetica", 11)
                _pdf_wrap(c, desc, x + 14, y - 78, cw - 28, 11, line_h=14)
        elif kind == "table":
            c.setFillColor(_hex(TEXT))
            c.setFont("Helvetica-Bold", 32)
            c.drawString(36, H - 100, slide["title"])
            y = H - 160
            col_w = (W - 72) / len(slide["headers"])
            c.setFillColor(_hex(SUB))
            c.setFont("Courier-Bold", 10)
            for ci, hdr in enumerate(slide["headers"]):
                c.drawString(36 + ci * col_w + 8, y, hdr)
            y -= 8
            c.setStrokeColor(_hex((40, 40, 40)))
            c.line(36, y, W - 36, y)
            y -= 24
            for row in slide["rows"]:
                for ci, val in enumerate(row):
                    if ci == 0:
                        c.setFillColor(_hex(TEXT))
                        c.setFont("Helvetica-Bold", 14)
                    else:
                        c.setFillColor(_hex(SUB))
                        c.setFont("Helvetica", 13)
                    c.drawString(36 + ci * col_w + 8, y, val)
                y -= 34
        elif kind == "ladder":
            c.setFillColor(_hex(TEXT))
            c.setFont("Helvetica-Bold", 32)
            c.drawString(36, H - 100, slide["title"])
            cw = (W - 72 - 30) / 3
            ch = 360
            for i, (tag, badge, body, color) in enumerate(slide["options"]):
                x = 36 + i * (cw + 15)
                y = H - 150
                c.setFillColor(_hex(CARD))
                c.rect(x, y - ch, cw, ch, fill=1, stroke=0)
                c.setFillColor(_hex(color))
                c.rect(x, y - ch, 5, ch, fill=1, stroke=0)
                c.setFont("Courier-Bold", 12)
                c.drawString(x + 18, y - 30, tag)
                c.setFillColor(_hex(SUB))
                c.setFont("Courier", 10)
                c.drawString(x + 18, y - 50, badge)
                c.setFillColor(_hex(TEXT))
                c.setFont("Helvetica", 13)
                _pdf_wrap(c, body, x + 18, y - 100, cw - 36, 13, line_h=18)
        elif kind == "weeks":
            c.setFillColor(_hex(TEXT))
            c.setFont("Helvetica-Bold", 32)
            c.drawString(36, H - 100, slide["title"])
            cw = (W - 72 - 30) / 3
            ch = 420
            for i, (tag, name, items) in enumerate(slide["weeks"]):
                x = 36 + i * (cw + 15)
                y = H - 150
                c.setFillColor(_hex(CARD))
                c.rect(x, y - ch, cw, ch, fill=1, stroke=0)
                c.setFillColor(_hex(GOLD))
                c.rect(x, y - 5, cw, 5, fill=1, stroke=0)
                c.setFillColor(_hex(GOLD))
                c.setFont("Courier-Bold", 12)
                c.drawString(x + 18, y - 30, tag)
                c.setFillColor(_hex(TEXT))
                c.setFont("Helvetica-Bold", 15)
                _pdf_wrap(c, name, x + 18, y - 60, cw - 36, 15, line_h=20)
                c.setFillColor(_hex(SUB))
                c.setFont("Helvetica", 12)
                ty = y - 130
                for it in items:
                    ty = _pdf_wrap(c, "• " + it, x + 18, ty, cw - 36, 12, line_h=18) - 4
        c.showPage()
    c.save()


# --------------------------------------------------------------------------
# Email .docx builder
# --------------------------------------------------------------------------
def build_email_docx(path: Path) -> None:
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = DocxPt(11)

    def _h(text, size=18, color=(20, 20, 20), bold=True, after=4):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.font.size = DocxPt(size)
        run.bold = bold
        run.font.color.rgb = DocxRGB(*color)
        p.paragraph_format.space_after = DocxPt(after)

    def _p(text, color=(40, 40, 40), size=11, after=8, bold=False):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.font.size = DocxPt(size)
        run.bold = bold
        run.font.color.rgb = DocxRGB(*color)
        p.paragraph_format.space_after = DocxPt(after)

    # Header block — meta
    _h("Email to Andrew Greene", size=20)
    _p("From: Lutar, Stephen P. · ORCID 0009-0001-0110-4173 · SZL Holdings", size=10, color=(110, 110, 110), after=2)
    _p("To: Andrew Greene, co-founder, Defense Unicorns", size=10, color=(110, 110, 110), after=2)
    _p("Date: Tuesday, 2026-05-19", size=10, color=(110, 110, 110), after=2)
    _p("Subject: a11oy.UDS — vision deck, architecture, and the meshing write-up I promised",
       size=10, color=(110, 110, 110), after=14)

    # Divider
    p = doc.add_paragraph()
    run = p.add_run("─" * 60)
    run.font.color.rgb = DocxRGB(200, 200, 200)
    p.paragraph_format.space_after = DocxPt(12)

    body_paragraphs = [
        "Andrew —",
        "Thanks for the time on the call last week and for the room you gave me to think this through end-to-end before sending. As promised on Friday, this email carries the a11oy.UDS package: a short vision deck, the architecture document, and the \"how I see it meshing in\" write-up.",
        "The recommendation is up front so you can decide whether the rest is worth your afternoon: start with Option A as a 2–3 week proof point, with Option C as the real destination. Option B falls out along the way. That ladder is the same one we sketched on the call.",
        "a11oy.UDS is the name I'd like us to use going forward — single token, native to UDS, inheriting your guardrails, carrying a11oy's orchestration DNA on top. The package is built to move the needle on two problems we've both named:",
    ]
    for para in body_paragraphs:
        _p(para, after=10)

    # Two-problem list
    p = doc.add_paragraph(style="List Number")
    run = p.add_run("Trusted AI/agent orchestration inside air-gapped UDS environments — provenance, human-in-the-loop approval gates, immutable tool-call audit, and disconnected operation, meshed with the UDS policy engine and Keycloak.")
    run.font.size = DocxPt(11)
    p = doc.add_paragraph(style="List Number")
    run = p.add_run("A UDS-native artifact spine for AI — SBOM-style attestation for models, prompts, embeddings, agent definitions, and evals; signed evals; drift detection; a promote / queue / discard flow that mirrors how Zarf already treats container images. The frontier-ingest + thesis-scoring layer inside a11oy is the working prototype.")
    run.font.size = DocxPt(11)

    _p(
        "The first-round proposal (the §00–§07 package you already have) and the Zarf wiring landed last sprint — uds-cli #5026 (in-bundle attestation), pepr #5027 (Λ-floor admission), and the three Zarf packages + UDS bundle under docs/proposals/defense-unicorns/szl-holdings/ are merged. The wires are set up. This Tuesday package is what gets built on top of them.",
        after=12,
    )

    _p("Attached (single zip — a11oy_uds_package.zip):", bold=True, after=4)
    for item in [
        "01_vision_deck.md — the deck outline (~14 slides) and the rendered .pptx / .pdf deck",
        "02_a11oy_uds_architecture.md — the architecture document, per-component table, problem-to-component map",
        "03_meshing_writeup.md — the meshing write-up (~1500 words)",
        "04_problem_briefs.md — one page per problem",
        "05_proof_plan.md — the 2–3 week Option A plan with week-by-week milestones",
        "06_appendix_evidence.md — the \"wires are set up\" exhibit list",
        "a11oy_uds_vision_deck.pptx + .pdf — the rendered deck",
    ]:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(item)
        run.font.size = DocxPt(11)

    _p(
        "Live view of the deck + architecture, no auth required, links back to all the source files: /uds inside a11oy.",
        after=12,
    )

    _p(
        "If A is welcome, I can have the proof-point payload demo-ready by Week 3 against any Mission App you point me at. If the answer is \"not yet\" or \"a different shape,\" I'd still value a 30-minute working session to recalibrate.",
        after=12,
    )

    _p("Talk soon,", after=2)
    _p("Stephen", after=10)

    _p("— Lutar, Stephen P. · ORCID 0009-0001-0110-4173 · SZL Holdings",
       size=10, color=(110, 110, 110))

    doc.save(path)


# --------------------------------------------------------------------------
# Zip builder
# --------------------------------------------------------------------------
ZIP_FILES = [
    "00_cover_letter.md",
    "01_vision_deck.md",
    "02_a11oy_uds_architecture.md",
    "03_meshing_writeup.md",
    "04_problem_briefs.md",
    "05_proof_plan.md",
    "06_appendix_evidence.md",
    "07_day_one_kickoff.md",
    "08_demo_script.md",
    "09_followup_responses.md",
    "a11oy_uds_vision_deck.pptx",
    "a11oy_uds_vision_deck.pdf",
    "README.txt",
]


README_TXT = """a11oy.UDS — Tuesday package for Andrew Greene
==============================================

From:    Lutar, Stephen P. · ORCID 0009-0001-0110-4173 · SZL Holdings
To:      Andrew Greene, co-founder, Defense Unicorns
Date:    Tuesday, 2026-05-19
Subject: a11oy.UDS — vision deck, architecture, meshing write-up

Read order
----------
  00_cover_letter.md            Email body Stephen pasted into the message.
  01_vision_deck.md             Vision deck outline (~14 slides + notes).
  02_a11oy_uds_architecture.md  Architecture: system view, per-component table,
                                problem-to-component map, AIArtifact CRD sketch.
  03_meshing_writeup.md         The meshing write-up (~1500 words). A / B / C
                                ladder, recommendation, 2–3 week proof plan.
  04_problem_briefs.md          One-page briefs for Problem 1 and Problem 2.
  05_proof_plan.md              Week-by-week proof plan for Option A.
  06_appendix_evidence.md       Evidence index: every "wires are set up" claim
                                points to a merged ref or a path on disk.

If Andrew says yes
------------------
  07_day_one_kickoff.md         The four asks for Defense Unicorns + the
                                first 48 hours, hour-by-hour.
  08_demo_script.md             Minute-by-minute Week-3 demo script + the
                                24-hour pre-flight checklist.
  09_followup_responses.md      Pre-drafted replies to Andrew's three most
                                likely responses (yes/yes+C/not yet) plus
                                the "acknowledged, will read later" nudge.

Rendered deck
-------------
  a11oy_uds_vision_deck.pptx    PowerPoint render of the vision deck.
  a11oy_uds_vision_deck.pdf     PDF render of the same deck (for offline view).

Live view
---------
  Inside the a11oy app at /uds — public, no auth, links back to every file
  above.

Recommendation
--------------
  Start with Option A as a 2–3 week proof point.
  Option C is the real destination.
  Option B falls out along the way.
"""


def build_zip(path: Path) -> None:
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for name in ZIP_FILES:
            fp = OUT / name
            if not fp.exists():
                continue
            zf.write(fp, arcname=f"a11oy_uds_package/{name}")


def main() -> int:
    OUT.mkdir(parents=True, exist_ok=True)
    (OUT / "README.txt").write_text(README_TXT, encoding="utf-8")
    print("Building vision deck (.pptx) …")
    build_pptx(OUT / "a11oy_uds_vision_deck.pptx")
    print("Building vision deck (.pdf) …")
    build_pdf(OUT / "a11oy_uds_vision_deck.pdf")
    print("Building email (.docx) …")
    build_email_docx(OUT / "email_to_andrew.docx")
    print("Building bundle zip …")
    build_zip(OUT / "a11oy_uds_package.zip")
    for f in OUT.iterdir():
        if f.is_file():
            print(f"  {f.relative_to(ROOT)}  ({f.stat().st_size:,} bytes)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
