"""Build the SZL Holdings APEX v2 Investor Pitch Deck as a polished .pptx.

Source: dossier/v2/APEX_v2_Investor_Pitch_Deck.md (16 slides + speaker notes)
Output: dossier/v2/pptx/APEX_v2_Investor_Pitch_Deck.pptx

Design: 16:9, dark navy background, gold accent rule, Helvetica/Calibri body.
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

NAVY      = RGBColor(0x0A, 0x16, 0x28)
NAVY_SOFT = RGBColor(0x12, 0x22, 0x38)
INK       = RGBColor(0xE8, 0xEC, 0xF1)
INK_DIM   = RGBColor(0x9A, 0xA6, 0xB8)
GOLD      = RGBColor(0xC9, 0xB7, 0x87)
GOLD_DARK = RGBColor(0xA8, 0x95, 0x60)
RULE      = RGBColor(0x2A, 0x38, 0x4E)
GREEN     = RGBColor(0x6B, 0xC4, 0x9A)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN_X = Inches(0.6)
TITLE_TOP = Inches(0.55)
EYEBROW_TOP = Inches(0.28)
BODY_TOP = Inches(1.85)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def add_bg(slide, color=NAVY):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg

def add_text(slide, x, y, w, h, text, *, size=18, color=INK, bold=False,
             font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w: s.line.width = line_w
    s.shadow.inherit = False
    return s

def add_footer(slide, page_no, total):
    # gold bottom rule
    add_rect(slide, MARGIN_X, Inches(7.05), SLIDE_W - 2*MARGIN_X, Emu(9525), fill=GOLD)
    add_text(slide, MARGIN_X, Inches(7.15), Inches(8), Inches(0.3),
             "SZL HOLDINGS  ·  GOVERNED OPERATIONAL INTELLIGENCE  ·  CONFIDENTIAL",
             size=8, color=INK_DIM, font="Calibri")
    add_text(slide, SLIDE_W - MARGIN_X - Inches(2.5), Inches(7.15), Inches(2.5), Inches(0.3),
             f"MAY 5, 2026   ·   {page_no:02d} / {total:02d}",
             size=8, color=INK_DIM, font="Calibri", align=PP_ALIGN.RIGHT)

def add_eyebrow(slide, text):
    add_text(slide, MARGIN_X, EYEBROW_TOP, Inches(8), Inches(0.3),
             text, size=10, color=GOLD, font="Calibri", bold=True)

def add_title(slide, title, subtitle=None):
    add_text(slide, MARGIN_X, TITLE_TOP, SLIDE_W - 2*MARGIN_X, Inches(0.9),
             title, size=32, color=INK, bold=True, font="Calibri")
    # gold underline
    add_rect(slide, MARGIN_X, Inches(1.5), Inches(0.6), Emu(28575), fill=GOLD)
    if subtitle:
        add_text(slide, MARGIN_X, Inches(1.55), SLIDE_W - 2*MARGIN_X, Inches(0.5),
                 subtitle, size=16, color=GOLD, font="Calibri", italic=True)

def add_notes(slide, note):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = note

def bullets(slide, x, y, w, h, items, *, size=16, color=INK, bullet_color=GOLD):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        p.space_after = Pt(8)
        # bullet marker as separate run
        r1 = p.add_run(); r1.text = "▸  "
        r1.font.name = "Calibri"; r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = bullet_color
        # parse **bold** segments
        parts = []
        rest = item
        while "**" in rest:
            pre, _, after = rest.partition("**")
            mid, _, post = after.partition("**")
            if pre: parts.append((pre, False))
            if mid: parts.append((mid, True))
            rest = post
        if rest: parts.append((rest, False))
        for txt, bld in parts:
            r = p.add_run(); r.text = txt
            r.font.name = "Calibri"; r.font.size = Pt(size); r.font.bold = bld
            r.font.color.rgb = INK if not bld else GOLD if bld else color
    return tb

def add_table(slide, x, y, w, h, headers, rows, *, header_fill=NAVY_SOFT,
              row_fill=None, alt_row_fill=None, font_size=12,
              first_col_bold=False, first_col_color=GOLD):
    n_rows = len(rows) + (1 if headers else 0)
    n_cols = len(headers) if headers else len(rows[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    # disable default theme stripes
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')
    if headers:
        for ci, head in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            tf = cell.text_frame; tf.clear()
            tf.margin_left = tf.margin_right = Inches(0.1)
            tf.margin_top = tf.margin_bottom = Inches(0.06)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = head
            r.font.name = "Calibri"; r.font.size = Pt(font_size); r.font.bold = True
            r.font.color.rgb = GOLD
    for ri, row in enumerate(rows):
        actual = ri + (1 if headers else 0)
        fill = alt_row_fill if (alt_row_fill and ri % 2 == 1) else row_fill
        for ci, val in enumerate(row):
            cell = tbl.cell(actual, ci)
            if fill is not None:
                cell.fill.solid(); cell.fill.fore_color.rgb = fill
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            tf = cell.text_frame; tf.clear()
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.1)
            tf.margin_top = tf.margin_bottom = Inches(0.08)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            # parse **bold**
            parts, rest = [], val
            while "**" in rest:
                pre, _, after = rest.partition("**")
                mid, _, post = after.partition("**")
                if pre: parts.append((pre, False))
                if mid: parts.append((mid, True))
                rest = post
            if rest: parts.append((rest, False))
            if not parts: parts = [(val, False)]
            bold_first = first_col_bold and ci == 0
            color_first = first_col_color if (first_col_bold and ci == 0) else INK
            for txt, bld in parts:
                r = p.add_run(); r.text = txt
                r.font.name = "Calibri"; r.font.size = Pt(font_size)
                r.font.bold = bld or bold_first
                r.font.color.rgb = color_first if bold_first else (GOLD if bld else INK)
    return tbl

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
TOTAL = 16

def slide1():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    # left gold bar
    add_rect(s, 0, 0, Inches(0.18), SLIDE_H, fill=GOLD)
    add_text(s, MARGIN_X+Inches(0.1), Inches(0.55), Inches(8), Inches(0.4),
             "APEX ACCELERATOR  ·  SERIES-A GRADE  ·  v2 REFRESH",
             size=11, color=GOLD, bold=True, font="Calibri")
    add_text(s, MARGIN_X+Inches(0.1), Inches(1.6), Inches(11.5), Inches(1.5),
             "SZL Holdings", size=72, color=INK, bold=True, font="Calibri")
    add_text(s, MARGIN_X+Inches(0.1), Inches(2.95), Inches(11.5), Inches(0.8),
             "Governed Operational Intelligence",
             size=30, color=GOLD, font="Calibri")
    add_text(s, MARGIN_X+Inches(0.1), Inches(3.85), Inches(11.5), Inches(1.0),
             "The platform that proves what AI decided,\nwhy, and on whose authority.",
             size=22, color=INK, italic=True, font="Calibri", line_spacing=1.25)
    # divider
    add_rect(s, MARGIN_X+Inches(0.1), Inches(5.5), Inches(2.0), Emu(19050), fill=GOLD)
    add_text(s, MARGIN_X+Inches(0.1), Inches(5.65), Inches(11.5), Inches(0.4),
             "Stephen Lutar  ·  Founder", size=16, color=INK, bold=True, font="Calibri")
    add_text(s, MARGIN_X+Inches(0.1), Inches(6.05), Inches(11.5), Inches(0.35),
             "stephenlutar2@gmail.com", size=12, color=INK_DIM, font="Calibri")
    add_text(s, MARGIN_X+Inches(0.1), Inches(6.35), Inches(11.5), Inches(0.35),
             "github.com/szl-holdings  ·  ORCID 0009-0001-0110-4173",
             size=12, color=INK_DIM, font="Calibri")
    add_text(s, MARGIN_X+Inches(0.1), Inches(6.65), Inches(11.5), Inches(0.35),
             "May 5, 2026", size=11, color=GOLD, bold=True, font="Calibri")
    add_notes(s, "Hold the slide for a beat. Single sentence: \"This is the platform that proves what AI decided, why, and on whose authority — by construction, not after the fact.\"")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Procurement Problem
# ─────────────────────────────────────────────────────────────────────────────
def slide2():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "01  ·  THE PROCUREMENT PROBLEM")
    add_title(s, "Every regulator is converging on the same question.",
              "\"Show your work.\"")
    bullets(s, MARGIN_X, Inches(2.7), SLIDE_W - 2*MARGIN_X, Inches(4),
            [
              "**NIST AI RMF**, **DoD RAI Strategy**, **OMB M-24-10**, **NY S.B. 7599** — all require runtime auditability at the point of decision.",
              "Existing AI vendors ship inference. They reconstruct audit trails after the fact.",
              "Human-in-the-loop is a checkbox in their docs. It is not a runtime gate.",
            ], size=18)
    # quote box
    add_rect(s, MARGIN_X, Inches(5.7), SLIDE_W - 2*MARGIN_X, Inches(1.0),
             fill=NAVY_SOFT, line=GOLD, line_w=Pt(0.5))
    add_text(s, MARGIN_X+Inches(0.3), Inches(5.85), SLIDE_W - 2*MARGIN_X - Inches(0.6), Inches(0.7),
             "The gap is not capability. The gap is proof — at the moment the decision is made.",
             size=14, color=GOLD, italic=True, font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 2, TOTAL)
    add_notes(s, "Land 'after the fact' hard. That is the gap.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — The Thesis
# ─────────────────────────────────────────────────────────────────────────────
def slide3():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "02  ·  THE THESIS")
    add_title(s, "The loop is the product.",
              "A bounded, governed, auditable decision loop — formalized as a mathematical invariant, shipped as a runtime.")
    # three-pillar layout
    pillars = [
        ("THE MATH", "Ouroboros Thesis v3, v9, v10",
         "DOI-pinned on Zenodo. CC-BY-4.0. Public, citable, reproducible."),
        ("THE RUNTIME", "@szl-holdings/ouroboros v6.2.0",
         "Open source. Full test suite. CodeQL clean. Live on GitHub."),
        ("THE AUDIT", "CPS Standard + Trust Plane + Agent Gateway",
         "The audit is the system itself — not a downstream report."),
    ]
    pw, gap = Inches(3.95), Inches(0.25)
    px = (SLIDE_W - 3*pw - 2*gap) / 2
    for i, (h, t, d) in enumerate(pillars):
        x = px + i*(pw + gap)
        add_rect(s, x, Inches(2.7), pw, Inches(3.6), fill=NAVY_SOFT, line=RULE, line_w=Pt(0.5))
        add_rect(s, x, Inches(2.7), pw, Emu(28575), fill=GOLD)
        add_text(s, x+Inches(0.3), Inches(2.95), pw-Inches(0.6), Inches(0.4),
                 h, size=11, color=GOLD, bold=True, font="Calibri")
        add_text(s, x+Inches(0.3), Inches(3.4), pw-Inches(0.6), Inches(1.0),
                 t, size=18, color=INK, bold=True, font="Calibri", line_spacing=1.2)
        add_text(s, x+Inches(0.3), Inches(4.7), pw-Inches(0.6), Inches(1.5),
                 d, size=13, color=INK_DIM, font="Calibri", line_spacing=1.35)
    add_footer(s, 3, TOTAL)
    add_notes(s, "This is the line that separates SZL from every other governed-AI pitch on the market. The math is written down.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Five Stages
# ─────────────────────────────────────────────────────────────────────────────
def slide4():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "03  ·  THE LOOP")
    add_title(s, "One loop. Five stages. Run continuously.")
    stages = [
        ("01", "INGEST", "Normalize signals across systems"),
        ("02", "SCORE", "Risk-score against policy in real time"),
        ("03", "DECIDE", "AI-supported decisions, policy-gated approvals"),
        ("04", "ACT", "Execute with full traceability — every input hashed, every approver named"),
        ("05", "VERIFY", "Verify outcomes against intent; replay-grade audit history"),
    ]
    cw, gap = Inches(2.35), Inches(0.12)
    cx = (SLIDE_W - 5*cw - 4*gap) / 2
    for i, (n, t, d) in enumerate(stages):
        x = cx + i*(cw+gap)
        add_rect(s, x, Inches(2.6), cw, Inches(3.7), fill=NAVY_SOFT, line=RULE, line_w=Pt(0.5))
        add_text(s, x+Inches(0.25), Inches(2.85), cw-Inches(0.5), Inches(0.6),
                 n, size=28, color=GOLD, bold=True, font="Calibri")
        add_text(s, x+Inches(0.25), Inches(3.55), cw-Inches(0.5), Inches(0.4),
                 t, size=14, color=INK, bold=True, font="Calibri")
        add_rect(s, x+Inches(0.25), Inches(4.0), Inches(0.4), Emu(19050), fill=GOLD)
        add_text(s, x+Inches(0.25), Inches(4.15), cw-Inches(0.5), Inches(2.0),
                 d, size=11, color=INK_DIM, font="Calibri", line_spacing=1.4)
        # arrow between
        if i < 4:
            ax = x + cw + Inches(0.005)
            add_text(s, ax, Inches(4.3), gap, Inches(0.4), "›",
                     size=22, color=GOLD, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(s, MARGIN_X, Inches(6.5), SLIDE_W-2*MARGIN_X, Inches(0.4),
             "Convergence properties formalized in the Lutar Invariant family — Ouroboros Thesis v9.",
             size=12, color=GOLD, italic=True, align=PP_ALIGN.CENTER, font="Calibri")
    add_footer(s, 4, TOTAL)
    add_notes(s, "The convergence properties of this loop — when it terminates, what its proof chain looks like — are formalized in the Lutar Invariant family in the published thesis.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Six Primitives
# ─────────────────────────────────────────────────────────────────────────────
def slide5():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "04  ·  THE PRIMITIVES")
    add_title(s, "Six primitives. Every vertical. Same shape.")
    prim = [
        ("Outcome Graph", "What was decided, by whom, with what evidence."),
        ("Proof Chain", "Hash chain from output → inputs → human approver."),
        ("Covenant Policy", "Declarative gating of every AI action."),
        ("Decision Simulation", "Replay decisions against history before policy ships."),
        ("Workflow Engine", "The bounded-loop scheduler."),
        ("Event Fabric (PRISM Bus)", "Append-only event spine. Everything publishes here."),
    ]
    cw, ch, gx, gy = Inches(3.95), Inches(2.0), Inches(0.2), Inches(0.2)
    cx = (SLIDE_W - 3*cw - 2*gx) / 2
    cy = Inches(2.55)
    for i, (t, d) in enumerate(prim):
        col = i % 3; row = i // 3
        x = cx + col*(cw+gx); y = cy + row*(ch+gy)
        add_rect(s, x, y, cw, ch, fill=NAVY_SOFT, line=RULE, line_w=Pt(0.5))
        add_rect(s, x, y, Emu(38100), ch, fill=GOLD)
        add_text(s, x+Inches(0.25), y+Inches(0.25), cw-Inches(0.5), Inches(0.5),
                 t, size=17, color=INK, bold=True, font="Calibri")
        add_text(s, x+Inches(0.25), y+Inches(0.85), cw-Inches(0.5), Inches(1.1),
                 d, size=12, color=INK_DIM, font="Calibri", line_spacing=1.4)
    add_footer(s, 5, TOTAL)
    add_notes(s, "This is what makes it a platform, not a portfolio of disconnected apps.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Eight Surfaces
# ─────────────────────────────────────────────────────────────────────────────
def slide6():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "05  ·  PRODUCT FOOTPRINT")
    add_title(s, "Eight live product surfaces. One control plane.")
    rows = [
        ["A11oy", "/a11oy/", "Orchestration · Decision Intelligence · Trust Plane"],
        ["Sentra (TENAX)", "/sentra/", "Cyber Resilience Command"],
        ["Amaru (Conduit)", "/conduit/", "Convergent Reverse-ETL"],
        ["Terra (DOMAINE)", "/terra/", "Real Estate Intelligence"],
        ["Vessels (SEXTANT)", "/vessels/", "Maritime Intelligence"],
        ["Counsel", "/counsel/", "Legal Matter Command"],
        ["Carlota Jo", "/carlota-jo/", "Concierge Advisory Operations"],
        ["ROSIE  ★ NEW", "/rosie/", "Unified Decision Fabric — operator surface for CPS payloads"],
    ]
    add_table(s, MARGIN_X, Inches(2.4), SLIDE_W-2*MARGIN_X, Inches(4.2),
              ["Surface", "Path", "Domain"], rows,
              header_fill=NAVY_SOFT, alt_row_fill=NAVY_SOFT, row_fill=NAVY,
              font_size=13, first_col_bold=True)
    add_footer(s, 6, TOTAL)
    add_notes(s, "Eight surfaces, not seven. ROSIE was promoted to a customer-facing artifact in the last 48 hours; it is where CPS payloads are operated.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — A11oy as Control Plane
# ─────────────────────────────────────────────────────────────────────────────
def slide7():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "06  ·  THE CONTROL PLANE")
    add_title(s, "A11oy is the Mythos-class control plane.")
    items = [
        ("ARGO",
         "Experience-era decision engine. Champion policies, self-play arena, mirror eval, counterfactuals, reward-hacking guardrails. World-model accuracy 89.1% · throughput 31.4 ev/s."),
        ("PSYCHE",
         "Emergent-sentience observatory."),
        ("TRUST CENTER · TRUST EXCHANGE · PUBLIC TRUST PORTAL",
         "Externally-facing proof distribution surface."),
        ("AGENT ZERO TRUST",
         "Live agent-gateway service enforcing OPA policy at the runtime boundary."),
        ("LIVE CLAUDE-CLASS ADVISOR  ·  /a11oy/chat",
         "System prompt locked to SOURCE_OF_TRUTH numbers; refuses fabricated metrics."),
    ]
    y = Inches(2.5)
    for h, d in items:
        add_rect(s, MARGIN_X, y, Emu(38100), Inches(0.85), fill=GOLD)
        add_text(s, MARGIN_X+Inches(0.25), y, Inches(11.5), Inches(0.4),
                 h, size=12, color=GOLD, bold=True, font="Calibri")
        add_text(s, MARGIN_X+Inches(0.25), y+Inches(0.4), Inches(11.5), Inches(0.5),
                 d, size=12, color=INK, font="Calibri", line_spacing=1.3)
        y += Inches(0.92)
    add_footer(s, 7, TOTAL)
    add_notes(s, "This is the slide that shows we are not just running models — we are governing the act of running them.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — CPS
# ─────────────────────────────────────────────────────────────────────────────
def slide8():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "07  ·  THE PROTOCOL")
    add_title(s, "Covenant Proof Standard (CPS)",
              "A payload-and-receipt protocol for cross-vertical governed actions.")
    api = [
        ("POST /api/cps/runs", "Execute a payload against a tenant"),
        ("POST /api/cps/runs/:id/approve", "Tier-checked human approval gate"),
        ("POST /api/cps/runs/:id/rollback", "Verified rollback to prior state"),
        ("POST /api/cps/payloads/:id/maturity", "Promote/demote payload maturity"),
    ]
    y = Inches(2.6)
    for ep, d in api:
        add_rect(s, MARGIN_X, y, SLIDE_W - 2*MARGIN_X, Inches(0.65),
                 fill=NAVY_SOFT, line=RULE, line_w=Pt(0.4))
        add_text(s, MARGIN_X+Inches(0.25), y+Inches(0.16), Inches(5.0), Inches(0.4),
                 ep, size=14, color=GOLD, bold=True, font="Consolas")
        add_text(s, MARGIN_X+Inches(5.5), y+Inches(0.18), Inches(7.5), Inches(0.4),
                 d, size=13, color=INK, font="Calibri")
        y += Inches(0.78)
    # callout
    add_rect(s, MARGIN_X, Inches(6.0), SLIDE_W - 2*MARGIN_X, Inches(0.85),
             fill=NAVY_SOFT, line=GOLD, line_w=Pt(0.75))
    add_text(s, MARGIN_X+Inches(0.3), Inches(6.05), SLIDE_W - 2*MARGIN_X - Inches(0.6), Inches(0.4),
             "THREE FLAGSHIP PAYLOADS LIVE TODAY",
             size=11, color=GOLD, bold=True, font="Calibri")
    add_text(s, MARGIN_X+Inches(0.3), Inches(6.40), SLIDE_W - 2*MARGIN_X - Inches(0.6), Inches(0.4),
             "Per-lane payloads rolling out across Vessels, Terra, Counsel, Carlota Jo.",
             size=13, color=INK, font="Calibri")
    add_footer(s, 8, TOTAL)
    add_notes(s, "CPS is what makes 'governed AI' a thing you can buy with a P.O. — not a thing you have to take on faith.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Three Papers
# ─────────────────────────────────────────────────────────────────────────────
def slide9():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "08  ·  THE MATH, WRITTEN DOWN")
    add_title(s, "Three peer-style papers. Public. DOI-pinned.")
    rows = [
        ["v3", "The Loop Is the Product",
         "Measuring Bounded Recursion as a System Primitive for Auditable AI",
         "Zenodo DOI 10.5281/zenodo.19944926  ·  CC-BY-4.0"],
        ["v9 ★", "The Lutar Invariant Family v1 → v7 → Ω",
         "From Three-Term Foundation to Bianchi-Closed Fiber Bundle",
         "17 pp · May 5 2026 · GitHub paper-v9-1.0.0 + Zenodo"],
        ["v10 ★", "The Audit Closure Operator Λ₁₀",
         "Formalising the Implementation Contract of the Lutar Family",
         "11 pp + Appendix A + Lutar one-pager · May 5 2026"],
    ]
    add_table(s, MARGIN_X, Inches(2.5), SLIDE_W-2*MARGIN_X, Inches(3.3),
              ["Paper", "Title", "Subtitle", "Status"], rows,
              header_fill=NAVY_SOFT, alt_row_fill=NAVY_SOFT, row_fill=NAVY,
              font_size=11, first_col_bold=True)
    # callout
    add_rect(s, MARGIN_X, Inches(6.0), SLIDE_W - 2*MARGIN_X, Inches(0.85),
             fill=NAVY_SOFT, line=GOLD, line_w=Pt(0.75))
    add_text(s, MARGIN_X+Inches(0.3), Inches(6.05), SLIDE_W - 2*MARGIN_X - Inches(0.6), Inches(0.8),
             "v10 is a meta-invariant on v9. It introduces no new physical term — its only job is to certify, layer by layer, that every formula in v9 actually executes against the live shipping repo.\nThe platform audits its own thesis.",
             size=12, color=INK, italic=True, font="Calibri", line_spacing=1.3)
    add_footer(s, 9, TOTAL)
    add_notes(s, "v10 is a meta-invariant on v9. It introduces no new physical term. Its only job is to certify, layer by layer, that every formula in v9 actually executes against the live shipping repo. The platform audits its own thesis.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Verified Numbers
# ─────────────────────────────────────────────────────────────────────────────
def slide10():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "09  ·  THE NUMBERS")
    add_title(s, "Re-runnable, not estimated.")
    metrics = [
        ("8", "Customer-facing surfaces", "+ A11oy orchestration"),
        ("6", "Platform primitives", "Same shape, every vertical"),
        ("848", "Database tables", "Provisioned & live"),
        ("5,524", "API endpoint declarations", "Across the monorepo"),
        ("126", "Monorepo packages", "Workspace total"),
        ("28", "Ouroboros packages", "Open-source runtime"),
        ("133", "Ouroboros runtime tests", "Live test calls"),
        ("62", "Ouroboros guardrails tests", "Bounded-loop invariants"),
        ("76 / 95", "Codex v11 nodes / edges", "11 domains"),
        ("126", "Security tests passing", "CodeQL clean"),
        ("23", "CI workflows", "Per-PR, per-merge"),
        ("3", "Peer-style papers", "DOI-pinned, CC-BY-4.0"),
    ]
    cw, ch, gx, gy = Inches(2.95), Inches(1.05), Inches(0.18), Inches(0.18)
    cx = (SLIDE_W - 4*cw - 3*gx) / 2
    cy = Inches(2.4)
    for i, (val, label, sub) in enumerate(metrics):
        col = i % 4; row = i // 4
        x = cx + col*(cw+gx); y = cy + row*(ch+gy)
        add_rect(s, x, y, cw, ch, fill=NAVY_SOFT, line=RULE, line_w=Pt(0.4))
        add_text(s, x+Inches(0.18), y+Inches(0.10), cw-Inches(0.36), Inches(0.45),
                 val, size=22, color=GOLD, bold=True, font="Calibri")
        add_text(s, x+Inches(0.18), y+Inches(0.55), cw-Inches(0.36), Inches(0.3),
                 label, size=11, color=INK, bold=True, font="Calibri")
        add_text(s, x+Inches(0.18), y+Inches(0.78), cw-Inches(0.36), Inches(0.25),
                 sub, size=9, color=INK_DIM, font="Calibri")
    add_text(s, MARGIN_X, Inches(6.55), SLIDE_W-2*MARGIN_X, Inches(0.4),
             "Every number ships with its verification command in SOURCE_OF_TRUTH.md.  We do not estimate.",
             size=12, color=GOLD, italic=True, align=PP_ALIGN.CENTER, font="Calibri")
    add_footer(s, 10, TOTAL)
    add_notes(s, "Every one of these numbers ships with the verification command in SOURCE_OF_TRUTH.md. We do not estimate.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Government Alignment
# ─────────────────────────────────────────────────────────────────────────────
def slide11():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "10  ·  GOVERNMENT ALIGNMENT")
    add_title(s, "Five demand vectors, mapped to shipped capability.")
    rows = [
        ["Operational transparency / IG audit",
         "Outcome Graph + Proof Chain + CPS payloads"],
        ["Cybersecurity / Zero Trust",
         "Sentra + governed adversary loop + agent gateway"],
        ["AI governance / OMB M-24-10 / NIST AI RMF",
         "Covenant Policy + Argo guardrails + mirror eval"],
        ["Cross-system data lineage",
         "Amaru + Frustum three-witness reconciliation"],
        ["Proof distribution to regulators",
         "Trust Center + Trust Exchange + Public Trust Portal"],
    ]
    add_table(s, MARGIN_X, Inches(2.5), SLIDE_W-2*MARGIN_X, Inches(3.6),
              ["Demand vector", "SZL surface"], rows,
              header_fill=NAVY_SOFT, alt_row_fill=NAVY_SOFT, row_fill=NAVY,
              font_size=14, first_col_bold=True, first_col_color=INK)
    add_text(s, MARGIN_X, Inches(6.3), SLIDE_W-2*MARGIN_X, Inches(0.5),
             "The Trust Plane is what closes the loop with the procurement officer — evidence packet, on demand, with chain-of-custody.",
             size=12, color=GOLD, italic=True, align=PP_ALIGN.CENTER, font="Calibri")
    add_footer(s, 11, TOTAL)
    add_notes(s, "The Trust Plane is what closes the loop with the procurement officer. It is the surface where they get the evidence packet, on demand, with chain-of-custody.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Honest Position
# ─────────────────────────────────────────────────────────────────────────────
def slide12():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "11  ·  HONEST POSITION")
    add_title(s, "What SZL Holdings is NOT claiming.")
    items = [
        "No active federal contracts.",
        "No federal cloud authorization, no ATO, no DoD impact-level designation.",
        "No external third-party audit.",
        "No signed platform customers.",
        "Single-founder operation.",
        "All numbers in this deck are internal platform metrics. Not revenue. Not users.",
    ]
    bullets(s, MARGIN_X, Inches(2.5), SLIDE_W-2*MARGIN_X, Inches(3.5), items, size=15)
    # bottom callout
    add_rect(s, MARGIN_X, Inches(5.85), SLIDE_W - 2*MARGIN_X, Inches(0.95),
             fill=NAVY_SOFT, line=GOLD, line_w=Pt(0.75))
    add_text(s, MARGIN_X+Inches(0.3), Inches(5.95), SLIDE_W - 2*MARGIN_X - Inches(0.6), Inches(0.8),
             "The strength of this position is the public proof.\nNot pretended traction.",
             size=16, color=GOLD, bold=True, italic=True, font="Calibri",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    add_footer(s, 12, TOTAL)
    add_notes(s, "This is the slide that earns trust with a procurement counselor in 30 seconds. Lead with what you don't have.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Demo Path
# ─────────────────────────────────────────────────────────────────────────────
def slide13():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "12  ·  DEMO PATH (5 MINUTES)")
    add_title(s, "One company. One design language. One proof spine.")
    rows = [
        ["0:00", "A11oy Trust Center", "Constitutional design — proof, covenants, attestation"],
        ["0:45", "A11oy Command Surface", "One operator pane, every domain"],
        ["1:30", "Amaru Dashboard", "Data fabric throughput (live Date.now() window)"],
        ["2:15", "Sentra Governed Adversary Loop", "Six-step proof chain Sentra ↔ A11oy"],
        ["3:00", "Counsel Matter Overview", "Vertical depth — legal command"],
        ["3:45", "Terra Distress Engine", "Vertical depth — real estate"],
        ["4:30", "Vessels Maritime Intelligence", "Vertical depth — maritime ops"],
    ]
    add_table(s, MARGIN_X, Inches(2.4), SLIDE_W-2*MARGIN_X, Inches(3.8),
              ["Min", "Surface", "Lands on"], rows,
              header_fill=NAVY_SOFT, alt_row_fill=NAVY_SOFT, row_fill=NAVY,
              font_size=12, first_col_bold=True)
    add_text(s, MARGIN_X, Inches(6.4), SLIDE_W-2*MARGIN_X, Inches(0.4),
             "Documented at docs/audits/INVESTOR_DEMO_PATH.md  ·  Per-artifact audits in docs/audits/",
             size=11, color=INK_DIM, italic=True, align=PP_ALIGN.CENTER, font="Calibri")
    add_footer(s, 13, TOTAL)
    add_notes(s, "Investors should ask 'show me.' This path runs the proof spine through five verticals in five minutes.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Roadmap
# ─────────────────────────────────────────────────────────────────────────────
def slide14():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "13  ·  90-DAY ROADMAP")
    add_title(s, "What ships next, on what calendar.")
    rows = [
        ["DAYS 1 – 14",
         "SAM.gov UEI activated. Primary NAICS confirmed."],
        ["DAYS 15 – 30",
         "CAGE issued. CMMC Level 1 self-assessment posted. First sources-sought scan pass."],
        ["DAYS 31 – 60",
         "First SBIR opportunity identified with APEX guidance. Per-lane CPS payloads complete for Vessels + Terra."],
        ["DAYS 61 – 90",
         "First sources-sought response submitted. CPS rollout complete. v11 thesis (planned) covering CPS as a procurement protocol."],
    ]
    y = Inches(2.45)
    for tag, txt in rows:
        add_rect(s, MARGIN_X, y, SLIDE_W - 2*MARGIN_X, Inches(1.0),
                 fill=NAVY_SOFT, line=RULE, line_w=Pt(0.4))
        add_rect(s, MARGIN_X, y, Emu(38100), Inches(1.0), fill=GOLD)
        add_text(s, MARGIN_X+Inches(0.3), y+Inches(0.2), Inches(2.5), Inches(0.4),
                 tag, size=12, color=GOLD, bold=True, font="Calibri")
        add_text(s, MARGIN_X+Inches(3.0), y+Inches(0.18), Inches(9.5), Inches(0.7),
                 txt, size=13, color=INK, font="Calibri", line_spacing=1.35,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.10)
    add_text(s, MARGIN_X, Inches(6.95), SLIDE_W-2*MARGIN_X, Inches(0.3),
             "If we miss a date, we say why and reset publicly.",
             size=11, color=GOLD, italic=True, align=PP_ALIGN.CENTER, font="Calibri")
    add_footer(s, 14, TOTAL)
    add_notes(s, "Calendar dates we will actually hit. If we do not hit a date, we say why and reset publicly.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — Ask
# ─────────────────────────────────────────────────────────────────────────────
def slide15():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_eyebrow(s, "14  ·  THE ASK")
    add_title(s, "What we need from you.")
    asks = [
        ("PROCUREMENT COUNSEL",
         "Most credible NYS pathway. Most credible federal pathway. Set-aside eligibility review."),
        ("CAPABILITY-STATEMENT FORMAT GUIDANCE",
         "What NYS and federal procurement officers actually want on the page."),
        ("TARGETING",
         "Which agencies buy this kind of work — and which buy it from sole-founder firms."),
        ("INTRODUCTIONS",
         "To the first procurement officer who will read the v9 thesis and the CPS API spec and ask the next question."),
    ]
    cw, ch, gx, gy = Inches(6.0), Inches(1.85), Inches(0.25), Inches(0.25)
    cx = (SLIDE_W - 2*cw - gx) / 2
    cy = Inches(2.45)
    for i, (h, d) in enumerate(asks):
        col = i % 2; row = i // 2
        x = cx + col*(cw+gx); y = cy + row*(ch+gy)
        add_rect(s, x, y, cw, ch, fill=NAVY_SOFT, line=RULE, line_w=Pt(0.5))
        add_rect(s, x, y, cw, Emu(28575), fill=GOLD)
        add_text(s, x+Inches(0.3), y+Inches(0.2), cw-Inches(0.6), Inches(0.4),
                 h, size=12, color=GOLD, bold=True, font="Calibri")
        add_text(s, x+Inches(0.3), y+Inches(0.65), cw-Inches(0.6), Inches(1.2),
                 d, size=13, color=INK, font="Calibri", line_spacing=1.4)
    add_text(s, MARGIN_X, Inches(6.7), SLIDE_W-2*MARGIN_X, Inches(0.4),
             "This is not a fundraising ask.  This is a sponsorship-of-readiness ask.",
             size=13, color=GOLD, italic=True, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
    add_footer(s, 15, TOTAL)
    add_notes(s, "This is not a fundraising ask. This is a sponsorship-of-readiness ask.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — Close
# ─────────────────────────────────────────────────────────────────────────────
def slide16():
    s = prs.slides.add_slide(blank); add_bg(s, NAVY)
    add_rect(s, 0, 0, Inches(0.18), SLIDE_H, fill=GOLD)
    add_text(s, MARGIN_X+Inches(0.1), Inches(0.55), Inches(8), Inches(0.4),
             "CLOSE  ·  16 / 16",
             size=11, color=GOLD, bold=True, font="Calibri")
    add_text(s, MARGIN_X+Inches(0.1), Inches(1.5), Inches(11.5), Inches(0.9),
             "The math is written down.",
             size=42, color=INK, bold=True, font="Calibri")
    add_text(s, MARGIN_X+Inches(0.1), Inches(2.45), Inches(11.5), Inches(0.9),
             "The runtime is open source.",
             size=42, color=INK, bold=True, font="Calibri")
    add_text(s, MARGIN_X+Inches(0.1), Inches(3.40), Inches(11.5), Inches(0.9),
             "The proof is a hash chain.",
             size=42, color=GOLD, bold=True, font="Calibri")
    add_rect(s, MARGIN_X+Inches(0.1), Inches(4.55), Inches(2.0), Emu(19050), fill=GOLD)
    add_text(s, MARGIN_X+Inches(0.1), Inches(4.7), Inches(11.5), Inches(0.7),
             "We are building the platform that makes governed AI procurable.",
             size=20, color=INK, italic=True, font="Calibri", line_spacing=1.25)
    # contact block
    cx = MARGIN_X+Inches(0.1)
    add_text(s, cx, Inches(5.85), Inches(11.5), Inches(0.35),
             "Stephen Lutar", size=14, color=INK, bold=True, font="Calibri")
    add_text(s, cx, Inches(6.15), Inches(11.5), Inches(0.3),
             "stephenlutar2@gmail.com  ·  github.com/szl-holdings  ·  ORCID 0009-0001-0110-4173",
             size=11, color=INK_DIM, font="Calibri")
    add_text(s, cx, Inches(6.50), Inches(11.5), Inches(0.3),
             "Ouroboros Runtime: github.com/szl-holdings/ouroboros (v6.2.0)",
             size=11, color=INK_DIM, font="Calibri")
    add_text(s, cx, Inches(6.80), Inches(11.5), Inches(0.3),
             "Ouroboros Thesis:  github.com/szl-holdings/ouroboros-thesis (v3, v9, v10)",
             size=11, color=INK_DIM, font="Calibri")
    add_notes(s, "Stop. Wait for questions. Don't keep talking.")

# Build all 16
for fn in [slide1, slide2, slide3, slide4, slide5, slide6, slide7, slide8,
           slide9, slide10, slide11, slide12, slide13, slide14, slide15, slide16]:
    fn()

os.makedirs("dossier/v2/pptx", exist_ok=True)
out = "dossier/v2/pptx/APEX_v2_Investor_Pitch_Deck.pptx"
prs.save(out)
print(f"Saved: {out}  ({os.path.getsize(out)} bytes, {len(prs.slides)} slides)")
